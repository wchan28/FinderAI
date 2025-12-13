from __future__ import annotations

from pathlib import Path
from typing import List, Dict, Union

from pptx import Presentation


def extract_text_from_pptx(file_path: Union[str, Path]) -> List[Dict]:
    """
    Extract text from a PowerPoint file.

    Returns a list of dicts, one per slide:
    [
        {"slide_number": 1, "text": "slide content..."},
        {"slide_number": 2, "text": "slide content..."},
    ]
    """
    file_path = Path(file_path)
    prs = Presentation(file_path)

    slides_content = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = "".join(run.text for run in paragraph.runs)
                    if paragraph_text.strip():
                        slide_text_parts.append(paragraph_text.strip())

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        slide_text_parts.append(" | ".join(row_text))

        slide_text = "\n".join(slide_text_parts)

        if slide_text.strip():
            slides_content.append({
                "slide_number": slide_index,
                "text": slide_text
            })

    return slides_content


def extract_full_text_from_pptx(file_path: Union[str, Path]) -> str:
    """
    Extract all text from a PowerPoint file as a single string.
    Includes slide markers for context.
    """
    slides = extract_text_from_pptx(file_path)

    full_text_parts = []
    for slide in slides:
        full_text_parts.append(f"[Slide {slide['slide_number']}]\n{slide['text']}")

    return "\n\n".join(full_text_parts)
